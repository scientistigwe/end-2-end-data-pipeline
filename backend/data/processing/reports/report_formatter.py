# backend/core/formatters/report_formatter.py

from __future__ import annotations
import logging
from typing import Dict, Any, Optional, List, Protocol, TypeVar, Callable, Union, Tuple
from datetime import datetime
from dataclasses import dataclass, field
from enum import Enum
import pandas as pd
import json
from pathlib import Path
from abc import ABC, abstractmethod

# Core imports
from backend.core.messaging.broker import MessageBroker
from backend.core.messaging.types import MessageType, ProcessingMessage

# Data pipeline imports
# Add this import
from backend.data_pipeline.analytics.analytics_processor import AnalyticsResult
from backend.data_pipeline.insight_analysis.insight_processor import InsightAnalysisResult

logger = logging.getLogger(__name__)

# Type variables for generic types
T = TypeVar('T')
FormatResult = TypeVar('FormatResult')


class ReportFormat(Enum):
    """Available report output formats"""
    JSON = "json"
    HTML = "html"
    PDF = "pdf"
    EXCEL = "excel"
    CSV = "csv"
    MARKDOWN = "markdown"
    XML = "xml"
    YAML = "yaml"
    TEXT = "text"
    POWERPOINT = "pptx"


class ReportType(Enum):
    """Types of reports that can be formatted"""
    # Analysis Reports
    QUALITY_ANALYSIS = "quality_analysis"
    INSIGHT_ANALYSIS = "insight_analysis"
    PERFORMANCE_ANALYSIS = "performance_analysis"

    # Summary Reports
    DATA_SUMMARY = "data_summary"
    PROCESS_SUMMARY = "process_summary"
    ERROR_SUMMARY = "error_summary"

    # Business Reports
    BUSINESS_INSIGHTS = "business_insights"
    KPI_REPORT = "kpi_report"
    TREND_ANALYSIS = "trend_analysis"

    # Technical Reports
    VALIDATION_REPORT = "validation_report"
    PROFILING_REPORT = "profiling_report"
    ANOMALY_REPORT = "anomaly_report"

    # Combined Reports
    COMPREHENSIVE = "comprehensive"
    EXECUTIVE_SUMMARY = "executive_summary"
    TECHNICAL_SUMMARY = "technical_summary"

    # Custom Reports
    CUSTOM = "custom"
    TEMPLATE_BASED = "template_based"


@dataclass
class ReportSection:
    """Structure for a report section"""
    title: str
    content: Any
    type: str = "text"
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ReportTemplate:
    """Structure for report templates"""
    name: str
    content: str
    format: ReportFormat
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class FormattingContext:
    """Context for report formatting process"""
    pipeline_id: str
    report_type: ReportType
    output_format: ReportFormat
    input_data: Dict[str, Any]
    template_name: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    formatted_output: Optional[Dict[str, Any]] = None
    sections: List[ReportSection] = field(default_factory=list)
    created_at: datetime = field(default_factory=datetime.now)
    updated_at: datetime = field(default_factory=datetime.now)
    status: str = "pending"
    error: Optional[str] = None


class FormatStrategy(Protocol[FormatResult]):
    """Protocol for format strategies"""

    def format(self, data: Dict[str, Any], template: Optional[str] = None) -> FormatResult:
        ...


class ReportFormatter:
    """
    Formats and presents analysis and insight results in various output formats
    """

    def __init__(self,
                 message_broker: MessageBroker,
                 template_dir: Union[str, Path] = "report_templates",
                 custom_formatters: Optional[Dict[str, FormatStrategy]] = None):
        """
        Initialize the report formatter

        Args:
            message_broker: Message broker for communication
            template_dir: Directory containing report templates
            custom_formatters: Optional dictionary of custom formatters
        """
        self.message_broker = message_broker
        self.logger = logging.getLogger(__name__)
        self.template_dir = Path(template_dir)
        self.custom_formatters = custom_formatters or {}

        # Track active formatting tasks
        self.active_tasks: Dict[str, FormattingContext] = {}

        # Initialize formatters
        self._initialize_formatters()

    def _initialize_formatters(self) -> None:
        """Initialize different output formatters"""
        self.formatters = {
            ReportFormat.JSON: self._format_json,
            ReportFormat.HTML: self._format_html,
            ReportFormat.PDF: self._format_pdf,
            ReportFormat.EXCEL: self._format_excel,
            ReportFormat.CSV: self._format_csv,
            ReportFormat.MARKDOWN: self._format_markdown,
            ReportFormat.XML: self._format_xml,
            ReportFormat.YAML: self._format_yaml,
            ReportFormat.TEXT: self._format_text,
            ReportFormat.POWERPOINT: self._format_powerpoint,
            **self.custom_formatters
        }

    def format_report(self,
                      pipeline_id: str,
                      report_type: str,
                      parameters: Dict[str, Any]) -> None:
        """
        Start report formatting process

        Args:
            pipeline_id: Pipeline identifier
            report_type: Type of report to generate
            parameters: Formatting parameters
        """
        try:
            # Validate report type and format
            try:
                report_type_enum = ReportType(report_type.upper())
                output_format = ReportFormat(parameters.get('format', 'json').upper())
            except KeyError as e:
                raise ValueError(f"Invalid report type or format: {str(e)}")

            # Create formatting context
            context = FormattingContext(
                pipeline_id=pipeline_id,
                report_type=report_type_enum,
                output_format=output_format,
                input_data=parameters.get('data', {}),
                template_name=parameters.get('template'),
                metadata=parameters.get('metadata', {})
            )

            self.active_tasks[pipeline_id] = context

            # Start formatting
            self._process_formatting(pipeline_id)

        except Exception as e:
            self.logger.error(f"Failed to start report formatting: {str(e)}")
            self._handle_formatting_error(pipeline_id, str(e))

    def _process_formatting(self, pipeline_id: str) -> None:
        """Process report formatting"""
        try:
            context = self.active_tasks[pipeline_id]

            # Prepare data based on report type
            prepared_data = self._prepare_data(context.report_type, context.input_data)

            # Get appropriate formatter
            formatter = self.formatters.get(context.output_format)
            if not formatter:
                raise ValueError(f"No formatter found for format: {context.output_format}")

            # Format report
            self._send_status_update(pipeline_id, "Formatting started")
            formatted_output = formatter(prepared_data, context.template_name)

            # Store results
            context.formatted_output = formatted_output
            context.status = "completed"
            context.updated_at = datetime.now()

            # Notify completion
            self._notify_completion(pipeline_id)

        except Exception as e:
            self._handle_formatting_error(pipeline_id, str(e))

    def _prepare_data(self,
                      report_type: ReportType,
                      input_data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare data based on report type"""
        preparation_methods = {
            ReportType.QUALITY_ANALYSIS: self._prepare_quality_analysis,
            ReportType.INSIGHT_ANALYSIS: self._prepare_insight_analysis,
            ReportType.BUSINESS_INSIGHTS: self._prepare_business_insights,
            ReportType.PERFORMANCE_ANALYSIS: self._prepare_performance_analysis,
            ReportType.COMPREHENSIVE: self._prepare_comprehensive_report,
            ReportType.EXECUTIVE_SUMMARY: self._prepare_executive_summary,
            ReportType.TECHNICAL_SUMMARY: self._prepare_technical_summary,
            ReportType.CUSTOM: lambda x: x  # Pass through for custom reports
        }

        prepare_method = preparation_methods.get(report_type)
        if not prepare_method:
            raise ValueError(f"No preparation method for report type: {report_type}")

        return prepare_method(input_data)

    # Format-specific methods
    def _format_json(self, data: Dict[str, Any],
                     template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as JSON"""
        return {
            'format': 'json',
            'content': json.dumps(data, indent=2),
            'generated_at': datetime.now().isoformat()
        }

    def _format_markdown(self, data: Dict[str, Any],
                         template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as Markdown"""
        template = self._load_template(template_name or 'default', 'md')
        markdown_content = self._render_markdown(template, data)
        return {
            'format': 'markdown',
            'content': markdown_content,
            'generated_at': datetime.now().isoformat()
        }

    def _format_yaml(self, data: Dict[str, Any],
                     template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as YAML"""
        try:
            import yaml
            return {
                'format': 'yaml',
                'content': yaml.dump(data, default_flow_style=False),
                'generated_at': datetime.now().isoformat()
            }
        except ImportError:
            raise ImportError("PyYAML is required for YAML formatting")

    def _render_markdown(self, template: str, data: Dict[str, Any]) -> str:
        """Render markdown template with data"""
        try:
            # Use a markdown templating library (e.g., jinja2)
            from jinja2 import Template
            template_obj = Template(template)
            return template_obj.render(**data)
        except ImportError:
            # Fallback to basic string formatting
            return template.format(**data)

    def add_custom_formatter(self, format_type: str,
                             formatter: FormatStrategy) -> None:
        """Add a custom formatter"""
        if not isinstance(format_type, str):
            raise ValueError("Format type must be a string")
        self.formatters[format_type] = formatter

    def remove_custom_formatter(self, format_type: str) -> None:
        """Remove a custom formatter"""
        if format_type in self.formatters:
            del self.formatters[format_type]

    def _prepare_analysis_summary(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare quality analysis summary"""
        return {
            'title': 'Data Quality Analysis Summary',
            'sections': [
                {
                    'title': 'Quality Metrics',
                    'content': data.get('quality_metrics', {})
                },
                {
                    'title': 'Issues Detected',
                    'content': data.get('issues', {})
                },
                {
                    'title': 'Recommendations',
                    'content': data.get('recommendations', [])
                }
            ]
        }

    def _prepare_insight_summary(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare business insight summary"""
        return {
            'title': 'Business Insights Summary',
            'sections': [
                {
                    'title': 'Key Insights',
                    'content': data.get('key_insights', [])
                },
                {
                    'title': 'Business Goal Analysis',
                    'content': data.get('business_insights', {})
                },
                {
                    'title': 'Additional Findings',
                    'content': data.get('additional_insights', {})
                }
            ]
        }

    def _prepare_combined_summary(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare combined summary"""
        return {
            'title': 'Combined Analysis Report',
            'quality_summary': self._prepare_analysis_summary(
                data.get('quality_results', {})
            ),
            'insight_summary': self._prepare_insight_summary(
                data.get('insight_results', {})
            )
        }

    def _format_html(self, data: Dict[str, Any],
                     template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as HTML"""
        template = self._load_template(template_name or 'default', 'html')
        rendered_html = self._render_template(template, data)

        return {
            'format': 'html',
            'content': rendered_html,
            'generated_at': datetime.now().isoformat()
        }

    def _format_pdf(self, data: Dict[str, Any],
                    template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as PDF"""
        # First generate HTML
        html_content = self._format_html(data, template_name)['content']
        # Convert to PDF (implementation depends on your PDF library choice)
        pdf_content = self._convert_to_pdf(html_content)

        return {
            'format': 'pdf',
            'content': pdf_content,
            'generated_at': datetime.now().isoformat()
        }

    def _format_excel(self, data: Dict[str, Any],
                      template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as Excel"""
        # Implementation depends on your Excel library choice
        return {
            'format': 'excel',
            'content': None,  # Add Excel content
            'generated_at': datetime.now().isoformat()
        }

    def _format_csv(self, data: Dict[str, Any],
                    template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as CSV"""
        try:
            # Use pandas to handle CSV conversion
            import pandas as pd
            import io

            # Convert the data to a DataFrame
            if isinstance(data, dict):
                # Handle nested dictionaries and lists
                flattened_data = self._flatten_dict(data)
                df = pd.DataFrame([flattened_data])
            elif isinstance(data, list):
                df = pd.DataFrame(data)
            else:
                df = pd.DataFrame(data)

            # Create a string buffer for CSV content
            csv_buffer = io.StringIO()
            df.to_csv(csv_buffer, index=False)
            csv_content = csv_buffer.getvalue()

            return {
                'format': 'csv',
                'content': csv_content,
                'generated_at': datetime.now().isoformat()
            }

        except Exception as e:
            self.logger.error(f"Error formatting CSV: {str(e)}")
            raise

    def _flatten_dict(self, d: Dict[str, Any], parent_key: str = '',
                      sep: str = '_') -> Dict[str, Any]:
        """Flatten nested dictionary for CSV conversion"""
        items: List[Tuple[str, Any]] = []
        for k, v in d.items():
            new_key = f"{parent_key}{sep}{k}" if parent_key else k
            if isinstance(v, dict):
                items.extend(self._flatten_dict(v, new_key, sep=sep).items())
            elif isinstance(v, list):
                # Convert list to string representation
                items.append((new_key, str(v)))
            else:
                items.append((new_key, v))
        return dict(items)

    def _format_powerpoint(self, data: Dict[str, Any],
                           template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import io

            # Create presentation
            prs = Presentation()

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = data.get('title', 'Report')

            # Content slides
            for section in data.get('sections', []):
                # Add new slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                # Set title
                slide.shapes.title.text = section.get('title', '')

                # Add content based on type
                content = section.get('content', {})
                if isinstance(content, dict):
                    self._add_dict_to_slide(slide, content)
                elif isinstance(content, list):
                    self._add_list_to_slide(slide, content)
                else:
                    self._add_text_to_slide(slide, str(content))

            # Save to bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)

            return {
                'format': 'pptx',
                'content': pptx_buffer.getvalue(),
                'generated_at': datetime.now().isoformat()
            }

        except ImportError:
            raise ImportError("python-pptx is required for PowerPoint formatting")
        except Exception as e:
            self.logger.error(f"Error formatting PowerPoint: {str(e)}")
            raise

    def _format_text(self, data: Dict[str, Any],
                     template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as plain text"""
        try:
            text_content = []

            # Add title
            title = data.get('title', 'Report')
            text_content.append(title)
            text_content.append('=' * len(title))
            text_content.append('')

            # Add sections
            for section in data.get('sections', []):
                # Section title
                section_title = section.get('title', '')
                text_content.append(section_title)
                text_content.append('-' * len(section_title))

                # Section content
                content = section.get('content', {})
                if isinstance(content, dict):
                    text_content.extend(self._dict_to_text(content))
                elif isinstance(content, list):
                    text_content.extend(self._list_to_text(content))
                else:
                    text_content.append(str(content))

                text_content.append('')  # Add spacing between sections

            return {
                'format': 'text',
                'content': '\n'.join(text_content),
                'generated_at': datetime.now().isoformat()
            }

        except Exception as e:
            self.logger.error(f"Error formatting text: {str(e)}")
            raise

    def _format_xml(self, data: Dict[str, Any],
                    template_name: Optional[str] = None) -> Dict[str, Any]:
        """Format report as XML"""
        try:
            from xml.etree.ElementTree import Element, SubElement, tostring
            from xml.dom import minidom

            def dict_to_xml(parent: Element, data: Dict[str, Any]) -> None:
                """Convert dictionary to XML elements"""
                for key, value in data.items():
                    child = SubElement(parent, self._sanitize_xml_tag(key))
                    if isinstance(value, dict):
                        dict_to_xml(child, value)
                    elif isinstance(value, list):
                        for item in value:
                            item_elem = SubElement(child, 'item')
                            if isinstance(item, dict):
                                dict_to_xml(item_elem, item)
                            else:
                                item_elem.text = str(item)
                    else:
                        child.text = str(value)

            # Create root element
            root = Element('report')

            # Add metadata
            metadata = SubElement(root, 'metadata')
            SubElement(metadata, 'generated_at').text = datetime.now().isoformat()
            SubElement(metadata, 'title').text = data.get('title', 'Report')

            # Add content
            content = SubElement(root, 'content')
            for section in data.get('sections', []):
                section_elem = SubElement(content, 'section')
                SubElement(section_elem, 'title').text = section.get('title', '')
                section_content = SubElement(section_elem, 'content')

                if isinstance(section.get('content'), dict):
                    dict_to_xml(section_content, section['content'])
                elif isinstance(section.get('content'), list):
                    items_elem = SubElement(section_content, 'items')
                    for item in section['content']:
                        item_elem = SubElement(items_elem, 'item')
                        if isinstance(item, dict):
                            dict_to_xml(item_elem, item)
                        else:
                            item_elem.text = str(item)
                else:
                    section_content.text = str(section.get('content', ''))

            # Convert to string with pretty printing
            xml_str = minidom.parseString(tostring(root)).toprettyxml(indent="  ")

            return {
                'format': 'xml',
                'content': xml_str,
                'generated_at': datetime.now().isoformat()
            }

        except Exception as e:
            self.logger.error(f"Error formatting XML: {str(e)}")
            raise

    def _sanitize_xml_tag(self, tag: str) -> str:
        """Sanitize string for use as XML tag"""
        # Replace invalid characters with underscores
        import re
        tag = re.sub(r'[^a-zA-Z0-9-_.]', '_', tag)
        # Ensure tag starts with a letter or underscore
        if tag[0].isdigit():
            tag = f'_{tag}'
        return tag

    # Helper methods for PowerPoint formatting
    def _add_dict_to_slide(self, slide: Any, data: Dict[str, Any],
                           left: float = 1, top: float = 2) -> None:
        """Add dictionary content to PowerPoint slide"""
        content = slide.placeholders[1].text_frame
        for key, value in data.items():
            p = content.add_paragraph()
            p.text = f"{key}: {value}"

    def _add_list_to_slide(self, slide: Any, data: List[Any],
                           left: float = 1, top: float = 2) -> None:
        """Add list content to PowerPoint slide"""
        content = slide.placeholders[1].text_frame
        for item in data:
            p = content.add_paragraph()
            p.text = f"• {item}"

    def _add_text_to_slide(self, slide: Any, text: str,
                           left: float = 1, top: float = 2) -> None:
        """Add text content to PowerPoint slide"""
        content = slide.placeholders[1].text_frame
        p = content.add_paragraph()
        p.text = text

    def _dict_to_text(self, data: Dict[str, Any], indent: int = 0) -> List[str]:
        """Convert dictionary to text lines"""
        lines = []
        for key, value in data.items():
            if isinstance(value, dict):
                lines.append('  ' * indent + f"{key}:")
                lines.extend(self._dict_to_text(value, indent + 1))
            else:
                lines.append('  ' * indent + f"{key}: {value}")
        return lines

    def _list_to_text(self, data: List[Any], indent: int = 0) -> List[str]:
        """Convert list to text lines"""
        return ['  ' * indent + f"- {item}" for item in data]

    def _load_template(self, template_name: str, format_type: str) -> str:
        """Load template file"""
        template_path = self.template_dir / f"{template_name}.{format_type}.template"
        if not template_path.exists():
            template_path = self.template_dir / f"default.{format_type}.template"

        return template_path.read_text()

    def _render_template(self, template: str, data: Dict[str, Any]) -> str:
        """Render template with data"""
        # Implement template rendering logic
        return template  # Placeholder

    def _convert_to_pdf(self, html_content: str) -> bytes:
        """Convert HTML to PDF"""
        # Implement PDF conversion logic
        return b""  # Placeholder

    def _send_status_update(self, pipeline_id: str, status: str) -> None:
        """Send status update message"""
        context = self.active_tasks.get(pipeline_id)
        if not context:
            return

        message = ProcessingMessage(
            message_type=MessageType.REPORT_UPDATE,
            content={
                'pipeline_id': pipeline_id,
                'report_type': context.report_type.value,
                'format': context.output_format.value,
                'status': status,
                'timestamp': datetime.now().isoformat()
            }
        )

        self.message_broker.publish(message)

    def _notify_completion(self, pipeline_id: str) -> None:
        """Notify completion of report formatting"""
        context = self.active_tasks.get(pipeline_id)
        if not context:
            return

        message = ProcessingMessage(
            message_type=MessageType.REPORT_COMPLETE,
            content={
                'pipeline_id': pipeline_id,
                'report_type': context.report_type.value,
                'format': context.output_format.value,
                'report_data': context.formatted_output,
                'metadata': context.metadata,
                'timestamp': datetime.now().isoformat()
            }
        )

        self.message_broker.publish(message)
        self._cleanup_task(pipeline_id)

    def _handle_formatting_error(self, pipeline_id: str, error: str) -> None:
        """Handle formatting errors"""
        message = ProcessingMessage(
            message_type=MessageType.REPORT_ERROR,
            content={
                'pipeline_id': pipeline_id,
                'error': error,
                'timestamp': datetime.now().isoformat()
            }
        )

        self.message_broker.publish(message)
        self._cleanup_task(pipeline_id)

    def _cleanup_task(self, pipeline_id: str) -> None:
        """Clean up formatting task resources"""
        if pipeline_id in self.active_tasks:
            del self.active_tasks[pipeline_id]

    def get_formatting_status(self, pipeline_id: str) -> Optional[Dict[str, Any]]:
        """Get current status of formatting task"""
        context = self.active_tasks.get(pipeline_id)
        if not context:
            return None

        return {
            'pipeline_id': pipeline_id,
            'report_type': context.report_type.value,
            'format': context.output_format.value,
            'status': context.status,
            'has_output': bool(context.formatted_output),
            'created_at': context.created_at.isoformat(),
            'updated_at': context.updated_at.isoformat()
        }

    def _prepare_quality_analysis(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare quality analysis report data"""
        return {
            'title': 'Quality Analysis Report',
            'sections': [
                {
                    'title': 'Analysis Results',
                    'content': data.get('results', {})
                },
                {
                    'title': 'Quality Metrics',
                    'content': data.get('metrics', {})
                },
                {
                    'title': 'Findings',
                    'content': data.get('findings', [])
                }
            ]
        }

    def _prepare_business_insights(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare business insights report data"""
        return {
            'title': 'Business Insights Report',
            'sections': [
                {
                    'title': 'Key Insights',
                    'content': data.get('insights', [])
                },
                {
                    'title': 'Metrics',
                    'content': data.get('metrics', {})
                },
                {
                    'title': 'Recommendations',
                    'content': data.get('recommendations', [])
                }
            ]
        }

    def _prepare_performance_analysis(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare performance analysis report data"""
        return {
            'title': 'Performance Analysis Report',
            'sections': [
                {
                    'title': 'Performance Metrics',
                    'content': data.get('metrics', {})
                },
                {
                    'title': 'Analysis Results',
                    'content': data.get('results', {})
                },
                {
                    'title': 'Recommendations',
                    'content': data.get('recommendations', [])
                }
            ]
        }

    def _prepare_comprehensive_report(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare comprehensive report data"""
        return {
            'title': 'Comprehensive Analysis Report',
            'sections': [
                {
                    'title': 'Executive Summary',
                    'content': data.get('summary', {})
                },
                {
                    'title': 'Quality Analysis',
                    'content': self._prepare_quality_analysis(data.get('quality', {}))
                },
                {
                    'title': 'Business Insights',
                    'content': self._prepare_business_insights(data.get('insights', {}))
                },
                {
                    'title': 'Performance Analysis',
                    'content': self._prepare_performance_analysis(data.get('performance', {}))
                }
            ]
        }

    def _prepare_executive_summary(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare executive summary report data"""
        return {
            'title': 'Executive Summary',
            'sections': [
                {
                    'title': 'Key Findings',
                    'content': data.get('findings', [])
                },
                {
                    'title': 'Highlights',
                    'content': data.get('highlights', {})
                },
                {
                    'title': 'Action Items',
                    'content': data.get('actions', [])
                }
            ]
        }

    def _prepare_technical_summary(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Prepare technical summary report data"""
        return {
            'title': 'Technical Summary',
            'sections': [
                {
                    'title': 'Technical Metrics',
                    'content': data.get('metrics', {})
                },
                {
                    'title': 'Technical Analysis',
                    'content': data.get('analysis', {})
                },
                {
                    'title': 'Technical Recommendations',
                    'content': data.get('recommendations', [])
                }
            ]
        }

    def __del__(self):
        """Cleanup formatter resources"""
        self.active_tasks.clear()

